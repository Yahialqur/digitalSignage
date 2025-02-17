from flask import Flask, request, jsonify, send_from_directory
import os
import shutil
import threading
import requests
from pptx import Presentation
from PIL import Image, ImageDraw
from io import BytesIO

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
SLIDES_FOLDER = "static/slides"
BRANCH_CLIENTS = [""]  # Add branch IPs

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(SLIDES_FOLDER, exist_ok=True)

@app.route('/upload', methods=['POST'])
def upload_pptx():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    
    pptx_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(pptx_path)
    
    # Clear existing slides
    shutil.rmtree(SLIDES_FOLDER)
    os.makedirs(SLIDES_FOLDER, exist_ok=True)
    
    convert_pptx_to_images(pptx_path)
    notify_branches()  # Trigger reload at all branches
    return jsonify({"message": "Slides uploaded and processed successfully"}), 200

def convert_pptx_to_images(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        img = slide_to_image(slide)
        img_path = os.path.join(SLIDES_FOLDER, f"slide_{i+1}.png")
        img.save(img_path, "PNG")

def slide_to_image(slide):
    width, height = 1280, 720  # Standard HD resolution
    img = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(img)

    # Extract text from slide
    text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
    
    # Draw text on image
    draw.text((50, 50), text, fill="black")

    return img

@app.route('/slides')
def get_slides():
    slides = sorted(os.listdir(SLIDES_FOLDER))
    return jsonify({"slides": slides})

@app.route('/slides/<filename>')
def serve_slide(filename):
    return send_from_directory(SLIDES_FOLDER, filename)

@app.route('/reload')
def reload_branches():
    notify_branches()
    return jsonify({"message": "Branches notified to reload"}), 200

def notify_branches():
    """Notify all branch clients to reload slides"""
    def notify():
        for branch in BRANCH_CLIENTS:
            try:
                requests.get(f"{branch}/reload")
            except requests.exceptions.RequestException:
                print(f"Failed to notify {branch}")

    threading.Thread(target=notify).start()

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
